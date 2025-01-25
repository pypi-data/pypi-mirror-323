import os
from pptx import Presentation
from pptx.util import Inches


class CreatePowerPoint:
    def __init__(self, presentation_title):
        self.figures = []
        self.presentation = Presentation()
        self.presentation_title = presentation_title

    def add_figure(self, figure):
        self.figures.append(figure)

    def create_presentation(self):
        output_dir = "./"
        os.makedirs(output_dir, exist_ok=True)
        # title slide
        title_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[0]
        )
        title = title_slide.shapes.title
        title.text = self.presentation_title

        for i, figure in enumerate(self.figures):
            # add blank slide
            slide = self.presentation.slides.add_slide(
                self.presentation.slide_layouts[5]
            )
            image_path = f"{output_dir}figure_{i+1}.png"
            # scale image to 3 for better image resolution
            figure.write_image(image_path, scale=3)
            slide.shapes.add_picture(
                image_path,
                Inches(0),
                Inches(0),
                Inches(10),
                Inches(7.5),
            )
            os.remove(image_path)
        # thank-you slide
        thank_you_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[0]
        )
        thank_you = thank_you_slide.shapes.title
        thank_you.text = "Thank You"

        # save the report
        output_file = os.path.join(output_dir, f"{self.presentation_title}.pptx")
        self.presentation.save(output_file)
        print(f"PowetPoint presentaion is created and saved to {output_dir}")
