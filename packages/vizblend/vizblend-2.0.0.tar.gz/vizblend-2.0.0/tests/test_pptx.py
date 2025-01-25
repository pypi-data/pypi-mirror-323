import os
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unittest.mock import MagicMock, patch
from vizblend.create_ppt_report import CreatePowerPoint


@pytest.fixture
def mock_figure():
    """Mock a Plotly figure with a write_image method."""
    figure = MagicMock()
    figure.write_image = MagicMock()
    return figure


@pytest.fixture
def create_presentation_instance():
    """Fixture for initializing CreatePowerPoint instance."""
    return CreatePowerPoint("Test Presentation")


def test_add_figure(create_presentation_instance, mock_figure):
    """Test that a figure is correctly added to the instance."""
    create_presentation_instance.add_figure(mock_figure)
    assert len(create_presentation_instance.figures) == 1
    assert create_presentation_instance.figures[0] == mock_figure


@patch("os.makedirs")
def test_create_presentation_handles_empty_figures(
    mock_makedirs, create_presentation_instance
):
    """Test that create_presentation works even when no figures are added."""
    with patch.object(create_presentation_instance.presentation, "save") as mock_save:
        create_presentation_instance.create_presentation()
        slides = create_presentation_instance.presentation.slides
        assert len(slides) == 2  # 1 title slide, 1 thank-you slide
        # Verify title slide text
        title_slide = slides[0]
        assert title_slide.shapes.title.text == "Test Presentation"
        # Verify thank-you slide text
        thank_you_slide = slides[-1]
        assert thank_you_slide.shapes.title.text == "Thank You"
        # Ensure presentation is saved
        mock_save.assert_called_once_with("./Test Presentation.pptx")


# @patch("os.makedirs")
# @patch("os.remove")
# @patch("os.path.join")
# def test_create_presentation(
#     mock_join, mock_remove, mock_makedirs, create_presentation_instance, mock_figure
# ):
#     """Test the creation of the PowerPoint presentation."""
#     # Add two mock figures
#     create_presentation_instance.add_figure(mock_figure)
#     create_presentation_instance.add_figure(mock_figure)

#     # Mock the write_image method to avoid actual file operations
#     with patch.object(mock_figure, "write_image", side_effect=lambda path: None):
#         # Mock os.path.join to return a specific output file path
#         mock_join.return_value = "./Test Presentation.pptx"

#         with patch.object(
#             create_presentation_instance.presentation, "save"
#         ) as mock_save:
#             create_presentation_instance.create_presentation()

#             # Check if title slide and thank-you slide are added
#             slides = create_presentation_instance.presentation.slides
#             assert len(slides) == 4  # 1 title slide, 2 figure slides, 1 thank-you slide

#             # Verify title slide text
#             title_slide = slides[0]
#             assert title_slide.shapes.title.text == "Test Presentation"

#             # Verify figure slide images are written and removed
#             assert mock_figure.write_image.call_count == 2
#             assert mock_remove.call_count == 2
#             mock_remove.assert_any_call("./figure_1.png")
#             mock_remove.assert_any_call("./figure_2.png")

#             # Verify thank-you slide text
#             thank_you_slide = slides[-1]
#             assert thank_you_slide.shapes.title.text == "Thank You"

#             # Ensure presentation is saved
#             mock_save.assert_called_once_with("./Test Presentation.pptx")


# @patch("os.makedirs")
# @patch("os.remove")
# @patch("os.path.join")
# def test_slide_content(
#     mock_join, mock_makedirs, mock_remove, create_presentation_instance, mock_figure
# ):
#     """Test that slides contain the expected elements."""
#     create_presentation_instance.add_figure(mock_figure)

#     # Mock the write_image method to avoid actual file operations
#     with patch.object(mock_figure, "write_image", side_effect=lambda path: None):
#         # Mock os.path.join to return a specific output file path
#         mock_join.return_value = "./Test Presentation.pptx"

#         with patch.object(create_presentation_instance.presentation, "save"):
#             create_presentation_instance.create_presentation()

#             slides = create_presentation_instance.presentation.slides

#             # Verify slide layout and content
#             figure_slide = slides[1]  # First figure slide
#             assert len(figure_slide.shapes) > 0
#             assert any(
#                 shape.shape_type == MSO_SHAPE_TYPE.PICTURE
#                 for shape in figure_slide.shapes
#             ), "Expected at least one picture on the slide"


@patch("os.makedirs")
def test_output_directory_creation(mock_makedirs, create_presentation_instance):
    """Test that the output directory is created if it does not exist."""
    with patch.object(create_presentation_instance.presentation, "save"):
        create_presentation_instance.create_presentation()
        mock_makedirs.assert_called_once_with("./", exist_ok=True)
